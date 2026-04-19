from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── palette ──────────────────────────────────────────────────────────
BG       = RGBColor(0x16, 0x16, 0x16)
CARD     = RGBColor(0x1f, 0x1f, 0x1f)
ACCENT   = RGBColor(0x3a, 0x7b, 0xd5)
GREEN    = RGBColor(0x7e, 0xc8, 0x7e)
WHITE    = RGBColor(0xff, 0xff, 0xff)
GREY     = RGBColor(0xaa, 0xaa, 0xaa)
DARK     = RGBColor(0x2d, 0x2d, 0x2d)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── helpers ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return txb


def bg(slide):
    add_rect(slide, 0, 0, W, H, BG)


def accent_bar(slide, y=Inches(0.55), h=Inches(0.06)):
    add_rect(slide, 0, y, W, h, ACCENT)


# ── slide 1 — title ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
add_rect(s, Inches(0), Inches(2.6), W, Inches(2.4), DARK)
add_text(s, "MBTI Persona Chat", Inches(1), Inches(2.7), Inches(11), Inches(1.2),
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "An AI-powered chat simulator using personality archetypes",
         Inches(1), Inches(3.7), Inches(11), Inches(0.6),
         size=20, color=GREY, align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.4), Inches(4.55), Inches(2.5), Inches(0.08), ACCENT)
add_text(s, "Built with Python · OpenAI GPT-4o · Flask",
         Inches(1), Inches(4.8), Inches(11), Inches(0.5),
         size=14, color=GREY, align=PP_ALIGN.CENTER, italic=True)


# ── slide 2 — what is it ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "What Is It?", Inches(0.7), Inches(0.7), Inches(10), Inches(0.7),
         size=34, bold=True, color=WHITE)

bullets = [
    ("🎭", "Persona Simulator", "Pick any of the 16 MBTI personality types and have a live AI-powered conversation that stays in character throughout."),
    ("💬", "Two Chat Modes",    "Chat as yourself against an AI persona, or watch two AI personas converse with each other automatically."),
    ("💾", "Full Logging",      "Every conversation is saved as a timestamped text file in a dedicated chats/ folder."),
    ("🌐", "Browser GUI",       "A clean dark-mode web interface served locally — no installs beyond Python packages."),
]

for i, (icon, title, body) in enumerate(bullets):
    x = Inches(0.6 + (i % 2) * 6.4)
    y = Inches(1.75 + (i // 2) * 2.4)
    add_rect(s, x, y, Inches(5.9), Inches(2.1), CARD)
    add_rect(s, x, y, Inches(0.12), Inches(2.1), ACCENT)
    add_text(s, icon + "  " + title, x + Inches(0.25), y + Inches(0.18),
             Inches(5.5), Inches(0.55), size=17, bold=True, color=WHITE)
    add_text(s, body, x + Inches(0.25), y + Inches(0.7),
             Inches(5.45), Inches(1.2), size=13, color=GREY)


# ── slide 3 — tech stack ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "Tech Stack", Inches(0.7), Inches(0.7), Inches(10), Inches(0.7),
         size=34, bold=True, color=WHITE)

stack = [
    ("OpenAI GPT-4o",   "Powers all persona responses with the chat completions API", ACCENT),
    ("Python 3.13",     "Core language for both CLI and server", GREEN),
    ("Flask",           "Lightweight web server serving the GUI and REST endpoints", RGBColor(0xff,0x7f,0x50)),
    ("python-dotenv",   "Loads OPENAI_API_KEY from a local .env file", RGBColor(0xf0,0xc0,0x30)),
    ("Vanilla JS + CSS","Frontend — no frameworks, dark-mode chat UI in pure HTML/JS", RGBColor(0xc0,0x80,0xff)),
]

for i, (name, desc, color) in enumerate(stack):
    y = Inches(1.55 + i * 1.02)
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.86), CARD)
    add_rect(s, Inches(0.6), y, Inches(0.12), Inches(0.86), color)
    add_text(s, name, Inches(0.9), y + Inches(0.08), Inches(3), Inches(0.4),
             size=16, bold=True, color=WHITE)
    add_text(s, desc, Inches(4.1), y + Inches(0.1), Inches(8.4), Inches(0.65),
             size=14, color=GREY)


# ── slide 4 — 16 personas ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "The 16 MBTI Personas", Inches(0.7), Inches(0.7), Inches(10), Inches(0.7),
         size=34, bold=True, color=WHITE)

types = [
    ("INTJ","Architect"), ("INTP","Thinker"),  ("ENTJ","Commander"), ("ENTP","Debater"),
    ("INFJ","Advocate"),  ("INFP","Mediator"),  ("ENFJ","Protagonist"),("ENFP","Campaigner"),
    ("ISTJ","Logistician"),("ISFJ","Defender"), ("ESTJ","Executive"), ("ESFJ","Consul"),
    ("ISTP","Virtuoso"),  ("ISFP","Adventurer"),("ESTP","Entrepreneur"),("ESFP","Entertainer"),
]

cols, rows = 4, 4
cw, rh = Inches(2.9), Inches(1.25)
for idx, (code, label) in enumerate(types):
    col = idx % cols
    row = idx // cols
    x = Inches(0.6) + col * (cw + Inches(0.22))
    y = Inches(1.6) + row * (rh + Inches(0.14))
    add_rect(s, x, y, cw, rh, CARD)
    add_rect(s, x, y, cw, Inches(0.06), ACCENT)
    add_text(s, code, x + Inches(0.15), y + Inches(0.12), cw - Inches(0.2), Inches(0.5),
             size=19, bold=True, color=ACCENT)
    add_text(s, label, x + Inches(0.15), y + Inches(0.58), cw - Inches(0.2), Inches(0.5),
             size=13, color=GREY)


# ── slide 5 — two modes ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "Two Chat Modes", Inches(0.7), Inches(0.7), Inches(10), Inches(0.7),
         size=34, bold=True, color=WHITE)

# mode 1
add_rect(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2), CARD)
add_rect(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.1), ACCENT)
add_text(s, "Human vs AI", Inches(0.85), Inches(1.75), Inches(5.3), Inches(0.6),
         size=22, bold=True, color=WHITE)
add_text(s,
    "• You pick a persona for yourself (or play as Human)\n"
    "• The AI plays the opposite persona\n"
    "• A random scenario is set at the start\n"
    "• The AI opens the conversation in character\n"
    "• You type freely — the AI stays in character\n"
    "• Full transcript saved to chats/",
    Inches(0.85), Inches(2.45), Inches(5.3), Inches(3.9),
    size=14, color=GREY)

# mode 2
add_rect(s, Inches(6.93), Inches(1.6), Inches(5.8), Inches(5.2), CARD)
add_rect(s, Inches(6.93), Inches(1.6), Inches(5.8), Inches(0.1), GREEN)
add_text(s, "AI vs AI (Auto)", Inches(7.18), Inches(1.75), Inches(5.3), Inches(0.6),
         size=22, bold=True, color=WHITE)
add_text(s,
    "• Pick two AI personas\n"
    "• Both personas get unique system prompts\n"
    "• They alternate turns automatically\n"
    "• 2.2 second delay between turns so you can read\n"
    "• Stop button pauses the conversation anytime\n"
    "• New Chat resets without closing the app",
    Inches(7.18), Inches(2.45), Inches(5.3), Inches(3.9),
    size=14, color=GREY)


# ── slide 6 — how it works ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "How It Works", Inches(0.7), Inches(0.7), Inches(10), Inches(0.7),
         size=34, bold=True, color=WHITE)

steps = [
    ("1", "Pick Personas",    "User selects two personas from the 16 MBTI types (or Human). Each AI persona gets a tailored system prompt that defines tone, speech style, and emotional register."),
    ("2", "Random Scenario",  "A scenario seed is randomly chosen (e.g. 'two strangers stuck at a flight delay'). Both AI prompts include this context to ground the conversation."),
    ("3", "Opening Line",     "The first AI persona calls GPT-4o to generate an opening message already mid-situation. This is injected as the first message in the chat history."),
    ("4", "Conversation Loop","Each turn appends to a shared history. For the active speaker, their own messages are 'assistant' and the other's are 'user'. This keeps both sides in character."),
]

for i, (num, title, body) in enumerate(steps):
    y = Inches(1.55 + i * 1.38)
    add_rect(s, Inches(0.6),  y, Inches(0.7), Inches(1.15), ACCENT)
    add_text(s, num, Inches(0.6), y + Inches(0.2), Inches(0.7), Inches(0.7),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(1.45), y, Inches(11.25), Inches(1.15), CARD)
    add_text(s, title, Inches(1.65), y + Inches(0.1), Inches(10.8), Inches(0.45),
             size=16, bold=True, color=WHITE)
    add_text(s, body, Inches(1.65), y + Inches(0.52), Inches(10.8), Inches(0.55),
             size=13, color=GREY)


# ── slide 7 — gui walkthrough ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
add_text(s, "GUI Walkthrough", Inches(0.7), Inches(0.7), Inches(10), Inches(0.7),
         size=34, bold=True, color=WHITE)

ui_parts = [
    ("Top Bar",        "Two persona dropdowns, a Start button, New Chat and Stop controls. Persona 1 sits on the left, Persona 2 on the right."),
    ("Scenario Banner","Green highlighted bar displaying the randomly chosen scenario — always visible so you know the context."),
    ("Chat Bubbles",   "Persona 1 messages appear on the left in blue, Persona 2 on the right in grey, and Human messages in green."),
    ("Typing Indicator","A '…' bubble appears while GPT-4o is generating, so the conversation feels live."),
    ("Input Bar",      "Shown only in Human mode. Press Enter or click Send. Disabled while waiting for AI reply."),
    ("Log Label",      "File path shown in the top bar so you always know where the transcript is being saved."),
]

for i, (part, desc) in enumerate(ui_parts):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.4)
    y = Inches(1.55 + row * 1.82)
    add_rect(s, x, y, Inches(5.9), Inches(1.6), CARD)
    add_rect(s, x, y, Inches(0.12), Inches(1.6), ACCENT)
    add_text(s, part, x + Inches(0.25), y + Inches(0.15), Inches(5.4), Inches(0.45),
             size=15, bold=True, color=WHITE)
    add_text(s, desc, x + Inches(0.25), y + Inches(0.6), Inches(5.4), Inches(0.9),
             size=13, color=GREY)


# ── slide 8 — summary ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s)
add_rect(s, 0, 0, W, H, BG)
add_rect(s, 0, Inches(2.8), W, Inches(2.0), DARK)
add_text(s, "What Was Built", Inches(1), Inches(2.9), Inches(11), Inches(0.8),
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

summary_items = [
    "CLI chat app  (main.py)",
    "Flask web GUI  (gui.py)",
    "16 MBTI personas + Human mode",
    "AI vs AI auto-chat with delay",
    "Automatic chat logging to chats/",
    "OpenAI GPT-4o · local .env key storage",
]
line = "     ·     ".join(summary_items[:3])
line2 = "     ·     ".join(summary_items[3:])
add_text(s, line,  Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.5),
         size=14, color=GREY, align=PP_ALIGN.CENTER)
add_text(s, line2, Inches(0.5), Inches(4.55), Inches(12.3), Inches(0.5),
         size=14, color=GREY, align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.4), Inches(5.25), Inches(2.5), Inches(0.08), ACCENT)


# ── save ────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "MBTI_Chat_Presentation.pptx")
prs.save(out)
print(f"saved: {out}")
